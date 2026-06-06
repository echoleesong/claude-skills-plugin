#!/usr/bin/env python3
"""
PPTX Generator Module for md-to-pptx

Generates PowerPoint presentations from structured slide data.
Supports: text, lists, code blocks, tables, images, and charts.
"""

import os
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass
from enum import Enum

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.chart.data import CategoryChartData
except ImportError:
    raise ImportError("python-pptx is required. Install with: pip install python-pptx")

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

# Import from md_parser (same directory)
from md_parser import SlideData, ContentElement, ContentType


class ThemeStyle(Enum):
    """Available theme styles."""
    BUSINESS = "business"
    TECH_DARK = "tech_dark"
    EDUCATION = "education"
    NEUMORPHISM = "neumorphism"


@dataclass
class ThemeConfig:
    """Theme configuration for styling."""
    name: str
    background_color: Tuple[int, int, int]
    title_color: Tuple[int, int, int]
    text_color: Tuple[int, int, int]
    accent_color: Tuple[int, int, int]
    code_bg_color: Tuple[int, int, int]
    title_font: str
    body_font: str
    title_size: int  # in points
    body_size: int
    code_size: int


# Built-in themes
THEMES: Dict[str, ThemeConfig] = {
    "business": ThemeConfig(
        name="Business",
        background_color=(255, 255, 255),
        title_color=(44, 62, 80),
        text_color=(52, 73, 94),
        accent_color=(41, 128, 185),
        code_bg_color=(236, 240, 241),
        title_font="Arial",
        body_font="Arial",
        title_size=32,
        body_size=14,
        code_size=12
    ),
    "tech_dark": ThemeConfig(
        name="Tech Dark",
        background_color=(30, 30, 30),
        title_color=(255, 255, 255),
        text_color=(220, 220, 220),
        accent_color=(0, 200, 150),
        code_bg_color=(45, 45, 45),
        title_font="Consolas",
        body_font="Segoe UI",
        title_size=32,
        body_size=14,
        code_size=12
    ),
    "education": ThemeConfig(
        name="Education",
        background_color=(255, 250, 240),
        title_color=(70, 130, 180),
        text_color=(60, 60, 60),
        accent_color=(255, 140, 0),
        code_bg_color=(245, 245, 245),
        title_font="Georgia",
        body_font="Verdana",
        title_size=32,
        body_size=14,
        code_size=12
    ),
    "neumorphism": ThemeConfig(
        name="Neumorphism",
        background_color=(240, 243, 249),
        title_color=(45, 55, 72),
        text_color=(74, 85, 104),
        accent_color=(66, 153, 225),
        code_bg_color=(226, 232, 240),
        title_font="Arial",
        body_font="Arial",
        title_size=32,
        body_size=14,
        code_size=12
    )
}


class PPTXGenerator:
    """Generate PPTX files from slide data."""

    # Slide dimensions (16:9 aspect ratio)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # Content margins
    MARGIN_LEFT = Inches(0.6)
    MARGIN_RIGHT = Inches(0.6)
    MARGIN_TOP = Inches(1.1)
    MARGIN_BOTTOM = Inches(0.4)

    def __init__(self, theme: str = "business", template_path: Optional[str] = None):
        """
        Initialize generator with theme or template.

        Args:
            theme: Theme name (business, tech_dark, education, neumorphism)
            template_path: Optional path to custom .pptx template
        """
        self.template_path = template_path
        self.theme = THEMES.get(theme, THEMES["business"])
        self.prs: Optional[Presentation] = None
        self.warnings: List[str] = []
        self.progress_callback = None
        self._use_template_background = template_path is not None

    def set_progress_callback(self, callback):
        """Set callback function for progress updates."""
        self.progress_callback = callback

    def _report_progress(self, current: int, total: int, message: str):
        """Report progress if callback is set."""
        if self.progress_callback:
            self.progress_callback(current, total, message)

    def generate(self, slides: List[SlideData], output_path: str) -> bool:
        """
        Generate PPTX from slide data.

        Args:
            slides: List of SlideData objects
            output_path: Output file path

        Returns:
            True if successful
        """
        self.warnings = []

        # Create presentation
        if self.template_path and os.path.exists(self.template_path):
            self.prs = Presentation(self.template_path)
            # Remove all existing template slides to start clean
            while len(self.prs.slides) > 0:
                rId = self.prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                self.prs.part.drop_rel(rId)
                self.prs.slides._sldIdLst.remove(self.prs.slides._sldIdLst[0])
        else:
            self.prs = Presentation()
            self.prs.slide_width = self.SLIDE_WIDTH
            self.prs.slide_height = self.SLIDE_HEIGHT

        total_slides = len(slides)

        for i, slide_data in enumerate(slides):
            self._report_progress(i + 1, total_slides, f"Generating slide {i + 1}: {slide_data.title[:30]}...")

            try:
                if slide_data.layout_hint == "title":
                    self._create_title_slide(slide_data)
                elif slide_data.layout_hint == "image":
                    self._create_image_slide(slide_data)
                elif slide_data.layout_hint == "chart":
                    self._create_chart_slide(slide_data)
                elif slide_data.layout_hint == "two_column":
                    self._create_two_column_slide(slide_data)
                else:
                    self._create_content_slide(slide_data)
            except Exception as e:
                self.warnings.append(f"Error on slide {i + 1}: {str(e)}")
                # Create fallback slide
                self._create_fallback_slide(slide_data, str(e))

        # Save presentation
        try:
            self.prs.save(output_path)
            self._report_progress(total_slides, total_slides, f"Saved to {output_path}")
            return True
        except Exception as e:
            self.warnings.append(f"Failed to save: {str(e)}")
            return False

    def _add_blank_slide(self) -> Any:
        """Add a blank slide to the presentation."""
        # Try to find a blank layout, fallback to last available layout
        layouts = self.prs.slide_layouts
        num_layouts = len(layouts)

        # Common blank layout indices to try
        blank_indices = [6, 5, num_layouts - 1, 0]

        for idx in blank_indices:
            if idx < num_layouts:
                try:
                    return self.prs.slides.add_slide(layouts[idx])
                except Exception:
                    continue

        # Last resort: use first available layout
        return self.prs.slides.add_slide(layouts[0])

    def _set_background(self, slide):
        """Set slide background color based on theme (skip if using template)."""
        if self._use_template_background:
            # Don't override template background
            return
        try:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*self.theme.background_color)
        except Exception:
            # Some templates may not support background modification
            pass

    def _add_title_shape(self, slide, title: str, top: float = 0.3) -> Any:
        """Add title text box to slide with dynamic height."""
        left = self.MARGIN_LEFT
        slide_w, _ = self._get_slide_dimensions()
        width = int(slide_w - self.MARGIN_LEFT - self.MARGIN_RIGHT)

        # Calculate height based on title length
        chars_per_line = max(30, int(width / Inches(1) * 4))
        est_lines = max(1, len(title) // chars_per_line + (1 if len(title) % chars_per_line else 0))
        height = max(0.6, est_lines * 0.45)

        shape = slide.shapes.add_textbox(left, Inches(top), width, Inches(height))
        tf = shape.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(self.theme.title_size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.theme.title_color)
        p.font.name = self.theme.title_font
        p.alignment = PP_ALIGN.LEFT

        return shape

    def _create_title_slide(self, slide_data: SlideData):
        """Create a title slide."""
        slide = self._add_blank_slide()
        self._set_background(slide)

        # Calculate title height based on text length
        title_text = slide_data.title
        title_font_size = 36  # Slightly smaller to fit better
        # Estimate lines: ~30 chars per line for centered title at 36pt
        chars_per_line = 35
        est_lines = max(1, len(title_text) // chars_per_line + (1 if len(title_text) % chars_per_line else 0))
        title_height = max(1.0, est_lines * 0.55)

        # Center title vertically - position it in upper third
        title_top = 2.0
        title_shape = slide.shapes.add_textbox(
            Inches(0.8),
            Inches(title_top),
            self.SLIDE_WIDTH - Inches(1.6),
            Inches(title_height)
        )
        tf = title_shape.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(title_font_size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.theme.title_color)
        p.font.name = self.theme.title_font
        p.alignment = PP_ALIGN.CENTER

        # Subtitle - positioned below title with clear gap
        if slide_data.subtitle:
            subtitle_top = title_top + title_height + 0.5
            subtitle_shape = slide.shapes.add_textbox(
                Inches(1.0),
                Inches(subtitle_top),
                self.SLIDE_WIDTH - Inches(2.0),
                Inches(1.0)
            )
            tf = subtitle_shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.subtitle
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(*self.theme.text_color)
            p.font.name = self.theme.body_font
            p.alignment = PP_ALIGN.CENTER

        # Add speaker notes
        if slide_data.notes:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = slide_data.notes

    def _create_content_slide(self, slide_data: SlideData):
        """Create a standard content slide with image-aware layout."""
        self._create_unified_slide(slide_data, use_two_column=False)

    def _get_slide_dimensions(self):
        """Get actual slide dimensions (accounts for template sizes)."""
        if self.prs:
            return self.prs.slide_width, self.prs.slide_height
        return self.SLIDE_WIDTH, self.SLIDE_HEIGHT

    def _get_image_bounds(self, element: ContentElement):
        """Get the absolute bounds (left, top, right, bottom) of a positioned image."""
        slide_w, slide_h = self._get_slide_dimensions()
        meta = element.metadata or {}
        position = meta.get("position", "")
        custom_width = self._parse_dimension(meta.get("width", ""), slide_w)
        custom_height = self._parse_dimension(meta.get("height", ""), slide_h)

        # Estimate image dimensions
        img_path = element.content
        target_width = custom_width if custom_width else Inches(5)
        target_height = custom_height if custom_height else Inches(4)

        if HAS_PIL and os.path.exists(img_path):
            try:
                with Image.open(img_path) as img:
                    img_w, img_h = img.size
                    aspect = img_w / img_h
                    if target_width / aspect > target_height:
                        final_h = target_height
                        final_w = target_height * aspect
                    else:
                        final_w = target_width
                        final_h = target_width / aspect
            except Exception:
                final_w, final_h = target_width, target_height
        else:
            final_w, final_h = target_width, target_height

        # Ensure integer EMU values
        final_w = int(round(final_w))
        final_h = int(round(final_h))

        left, top = self._resolve_image_position(position, final_w, final_h, meta)
        left = int(round(left))
        top = int(round(top))
        return left, top, left + final_w, top + final_h

    def _get_image_h_bounds(self, images: list):
        """Get the combined horizontal bounds of all positioned images."""
        if not images:
            return self.MARGIN_LEFT, self.SLIDE_WIDTH - self.MARGIN_RIGHT
        min_left = self.SLIDE_WIDTH
        max_right = 0
        for img_elem in images:
            left, _, right, _ = self._get_image_bounds(img_elem)
            min_left = min(min_left, left)
            max_right = max(max_right, right)
        return min_left, max_right

    def _calc_content_area(self, positioned_images, content_top, default_left, default_width, slide_bottom, img_gap):
        """
        Calculate available content area that avoids positioned images (using metadata estimates).

        Returns (content_left, content_width, adjusted_content_top).
        """
        if not positioned_images:
            return default_left, default_width, content_top

        default_right = default_left + default_width

        for img_elem in positioned_images:
            img_left, img_top, img_right, img_bottom = self._get_image_bounds(img_elem)

            # Check if image overlaps vertically with content area
            if img_bottom <= content_top or img_top >= slide_bottom:
                continue

            position = img_elem.metadata.get("position", "")
            if position in ("top-center", "center", "bottom-center"):
                new_top = img_bottom + img_gap
                if new_top > content_top:
                    content_top = new_top
            elif position in ("center-right", "top-right", "bottom-right"):
                new_right = img_left - img_gap
                if new_right > default_left + default_width * 0.3:
                    default_width = new_right - default_left
            elif position in ("center-left", "top-left", "bottom-left"):
                new_left = img_right + img_gap
                if new_left < default_right - default_width * 0.3:
                    default_width = default_right - new_left
                    default_left = new_left

        return default_left, default_width, content_top

    def _calc_content_area_from_bounds(self, img_bounds_list, content_top, default_left, default_width, slide_bottom, img_gap):
        """
        Calculate available content area using actual placed image bounds.

        Args:
            img_bounds_list: List of dicts with 'left', 'top', 'right', 'bottom', 'position' keys

        Returns (content_left, content_width, adjusted_content_top).
        """
        if not img_bounds_list:
            return default_left, default_width, content_top

        default_right = default_left + default_width

        for bounds in img_bounds_list:
            il = bounds['left']
            it = bounds['top']
            ir = bounds['right']
            ib = bounds['bottom']
            position = bounds['position']

            # Check if image overlaps vertically with content area
            if ib <= content_top or it >= slide_bottom:
                continue

            if position in ("top-center", "center", "bottom-center"):
                new_top = ib + img_gap
                if new_top > content_top:
                    content_top = new_top
            elif position in ("center-right", "top-right", "bottom-right"):
                new_right = il - img_gap
                if new_right > default_left + default_width * 0.3:
                    default_width = int(new_right - default_left)
            elif position in ("center-left", "top-left", "bottom-left"):
                new_left = ir + img_gap
                if new_left < default_right - default_width * 0.3:
                    default_width = int(default_right - new_left)
                    default_left = int(new_left)

        return int(default_left), int(default_width), int(content_top)

    def _create_image_slide(self, slide_data: SlideData):
        """Create a slide with image focus."""
        self._create_unified_slide(slide_data, use_two_column=False)

    def _create_chart_slide(self, slide_data: SlideData):
        """Create a slide with chart focus."""
        slide = self._add_blank_slide()
        self._set_background(slide)

        # Add title
        self._add_title_shape(slide, slide_data.title)

        # Find chart or table data
        chart_elem = None
        for elem in slide_data.elements:
            if elem.type in (ContentType.CHART_DATA, ContentType.TABLE):
                chart_elem = elem
                break

        if chart_elem:
            if chart_elem.type == ContentType.TABLE:
                # Convert table to chart
                self._add_chart_from_table(
                    slide, chart_elem,
                    Inches(1), Inches(1.5),
                    self.SLIDE_WIDTH - Inches(2),
                    Inches(5)
                )
            else:
                self._add_chart_element(
                    slide, chart_elem,
                    Inches(1), Inches(1.5),
                    self.SLIDE_WIDTH - Inches(2),
                    Inches(5)
                )

        # Add speaker notes
        if slide_data.notes:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = slide_data.notes

    def _create_two_column_slide(self, slide_data: SlideData):
        """Create a two-column layout slide."""
        self._create_unified_slide(slide_data, use_two_column=True)

    def _create_unified_slide(self, slide_data: SlideData, use_two_column: bool = False):
        """Unified slide layout engine.

        Handles all content layouts with proper image-text integration:
        - Center/top-center images: placed at top, text flows below at full width
        - Top-right images: text on left, expands to full width after image bottom
        - Two-column: logical element grouping with balanced height split
        - Per-element overflow detection
        """
        slide = self._add_blank_slide()
        self._set_background(slide)

        slide_w, slide_h = self._get_slide_dimensions()

        # Title
        title_shape = self._add_title_shape(slide, slide_data.title)
        title_bottom = int(Inches(0.3) + title_shape.height + Inches(0.2))

        content_top = title_bottom
        img_gap = int(Inches(0.2))

        # Separate positioned images from flow elements
        pos_images = []  # (element, index)
        flow_elements = []
        for idx, elem in enumerate(slide_data.elements):
            if elem.type == ContentType.IMAGE and elem.metadata.get("position"):
                pos_images.append((elem, idx))
            else:
                flow_elements.append(elem)

        # Place positioned images and collect actual bounds
        placed_img_bounds = []
        for elem, _ in pos_images:
            meta = elem.metadata or {}
            img_path = elem.content
            if not os.path.exists(img_path):
                self.warnings.append(f"Image not found: {img_path}")
                continue

            position = meta.get("position", "")
            custom_w = self._parse_dimension(meta.get("width", ""), slide_w)
            custom_h = self._parse_dimension(meta.get("height", ""), slide_h)
            target_w = custom_w if custom_w else int(Inches(5))
            max_h = custom_h if custom_h else int(Inches(4))

            # Compute image size preserving aspect ratio
            if HAS_PIL:
                try:
                    with Image.open(img_path) as pil_img:
                        iw, ih = pil_img.size
                        aspect = iw / ih
                        if target_w / aspect > max_h:
                            final_h = int(round(max_h))
                            final_w = int(round(max_h * aspect))
                        else:
                            final_w = int(round(target_w))
                            final_h = int(round(target_w / aspect))
                except Exception:
                    final_w, final_h = int(round(target_w)), int(round(max_h))
            else:
                final_w, final_h = int(round(target_w)), int(round(max_h))

            # Position with title-aware offset
            img_left, img_top = self._resolve_image_position(
                position, final_w, final_h, meta, title_top_offset=title_bottom
            )
            img_left = int(round(img_left))
            img_top = int(round(img_top))

            slide.shapes.add_picture(img_path, img_left, img_top, final_w, final_h)
            placed_img_bounds.append({
                'left': img_left, 'top': img_top,
                'right': img_left + final_w, 'bottom': img_top + final_h,
                'position': position,
            })

        # Build content zones from image bounds
        default_left = int(self.MARGIN_LEFT)
        default_width = int(slide_w - self.MARGIN_LEFT - self.MARGIN_RIGHT)
        slide_bottom = int(slide_h - self.MARGIN_BOTTOM)

        # Sort by top for consistent zone processing
        placed_img_bounds.sort(key=lambda b: b['top'])

        zones = []
        cur_top = content_top
        cur_left = default_left
        cur_width = default_width

        for bounds in placed_img_bounds:
            il, it, ir, ib = bounds['left'], bounds['top'], bounds['right'], bounds['bottom']
            pos = bounds['position']

            if pos in ("center", "top-center"):
                # Image blocks full width → ALL content goes below
                zone_bottom = ib + img_gap
                if zone_bottom > cur_top:
                    cur_top = zone_bottom
                    cur_left = default_left
                    cur_width = default_width
            elif pos in ("top-right", "center-right"):
                # Left zone: content beside image
                if it > cur_top + int(Inches(0.3)):
                    zones.append({'top': cur_top, 'left': cur_left,
                                  'width': cur_width, 'bottom': it})
                # Narrow zone while image is present
                avail_w = il - img_gap - cur_left
                if avail_w > default_width * 0.25:
                    zones.append({'top': max(cur_top, it), 'left': cur_left,
                                  'width': int(avail_w), 'bottom': ib + img_gap})
                cur_top = ib + img_gap
                cur_left = default_left
                cur_width = default_width
            elif pos in ("top-left", "center-left"):
                if it > cur_top + int(Inches(0.3)):
                    zones.append({'top': cur_top, 'left': cur_left,
                                  'width': cur_width, 'bottom': it})
                avail_l = ir + img_gap
                avail_w = default_left + default_width - avail_l
                if avail_w > default_width * 0.25:
                    zones.append({'top': max(cur_top, it), 'left': int(avail_l),
                                  'width': int(avail_w), 'bottom': ib + img_gap})
                cur_top = ib + img_gap
                cur_left = default_left
                cur_width = default_width

        # Final zone (remaining space)
        if cur_top < slide_bottom:
            zones.append({'top': cur_top, 'left': cur_left,
                          'width': cur_width, 'bottom': slide_bottom})

        if not zones:
            zones.append({'top': content_top, 'left': default_left,
                          'width': default_width, 'bottom': slide_bottom})

        # --- Place flow elements into zones ---
        if use_two_column and len(flow_elements) >= 3:
            self._place_two_column(slide, flow_elements, zones, slide_bottom)
        else:
            self._place_sequential(slide, flow_elements, zones, slide_bottom)

        # Speaker notes
        if slide_data.notes:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = slide_data.notes

    def _estimate_element_height(self, element: ContentElement, width: int) -> int:
        """Estimate the height an element will consume when placed."""
        if element.type == ContentType.HEADING:
            return int(Inches(0.45))
        elif element.type == ContentType.PARAGRAPH:
            text = element.content
            chars_per_line = max(40, int(width / Inches(1) * 5))
            est_lines = max(1, len(text) // chars_per_line + (1 if len(text) % chars_per_line else 0))
            return int(Inches(max(0.3, est_lines * 0.28)))
        elif element.type in (ContentType.BULLET_LIST, ContentType.NUMBERED_LIST):
            items = element.content
            chars_per_line = max(40, int(width / Inches(1) * 5))
            total_lines = 0
            for item in items:
                text_len = len(item['text']) + item.get("indent", 0) * 4 + 3
                total_lines += max(1, text_len // chars_per_line + (1 if text_len % chars_per_line else 0))
            return int(Inches(0.3) * total_lines)
        elif element.type == ContentType.TABLE:
            data = element.content
            rows = data.get("rows", [])
            headers = data.get("headers", [])
            row_count = len(rows) + (1 if headers else 0)
            return int(Inches(0.33) * row_count)
        elif element.type == ContentType.CODE_BLOCK:
            lines = element.content.split('\n')
            line_count = min(len(lines), 15)
            return int(Inches(0.25) * line_count + Inches(0.3))
        elif element.type == ContentType.IMAGE:
            return int(Inches(3))
        return int(Inches(0.3))

    def _find_zone_for_y(self, cur_y: int, zones: list) -> dict:
        """Find the zone that contains cur_y, or the next available zone."""
        gap = int(Inches(0.15))
        # First: find zone that contains cur_y (top <= cur_y < bottom)
        for z in zones:
            if z['top'] <= cur_y < z['bottom']:
                return z
        # Second: find the next zone whose top is near or after cur_y
        for z in zones:
            if z['top'] >= cur_y - gap:
                return z
        # Fallback: last zone
        return zones[-1] if zones else None

    def _place_sequential(self, slide, elements: list, zones: list, slide_bottom: int):
        """Place elements sequentially through content zones."""
        if not zones:
            return
        cur_y = zones[0]['top']

        for elem in elements:
            zone = self._find_zone_for_y(cur_y, zones)
            if zone is None:
                break

            # If we jumped to a new zone, reset cur_y to zone top
            if cur_y < zone['top']:
                cur_y = zone['top']

            # Check overflow
            avail_h = zone['bottom'] - cur_y
            est_h = self._estimate_element_height(elem, zone['width'])
            if avail_h < int(Inches(0.3)) or (est_h > avail_h and avail_h < int(Inches(0.8))):
                # Try next zone
                try:
                    zone_idx = zones.index(zone)
                except ValueError:
                    zone_idx = -1
                if zone_idx + 1 < len(zones):
                    zone = zones[zone_idx + 1]
                    cur_y = zone['top']
                    avail_h = zone['bottom'] - cur_y
                    if avail_h < int(Inches(0.3)):
                        break
                else:
                    break

            max_h = int(max(Inches(0.3), avail_h - Inches(0.1)))
            h = self._add_element(slide, elem, zone['left'], cur_y, zone['width'], max_h)

            if h > 0:
                cur_y += int(h) + int(Inches(0.08))

    def _place_two_column(self, slide, elements: list, zones: list, slide_bottom: int):
        """Place elements in two-column layout with logical grouping."""
        zone = zones[0]
        col_top = zone['top']
        col_width = int((zone['width'] - Inches(0.4)) / 2)
        col_left_1 = zone['left']
        col_left_2 = zone['left'] + col_width + int(Inches(0.4))

        # Group elements: heading + following non-heading = one unit
        groups = []  # List of (elements_list, weight)
        cur_group = []
        cur_weight = 0

        for elem in elements:
            if elem.type == ContentType.HEADING and cur_group:
                groups.append((cur_group, cur_weight))
                cur_group = [elem]
                cur_weight = self._estimate_element_height(elem, col_width)
            else:
                cur_group.append(elem)
                cur_weight += self._estimate_element_height(elem, col_width)
        if cur_group:
            groups.append((cur_group, cur_weight))

        # Split groups into two balanced columns
        total_w = sum(w for _, w in groups)
        target_w = total_w / 2

        col1_groups = []
        col2_groups = []
        cum_w = 0
        split_done = False

        for i, (grp, w) in enumerate(groups):
            if not split_done:
                col1_groups.append(grp)
                cum_w += w
                if cum_w >= target_w and i < len(groups) - 1:
                    split_done = True
            else:
                col2_groups.append(grp)

        # If only one column has content, fall back to sequential
        if not col2_groups:
            self._place_sequential(slide, elements, zones, slide_bottom)
            return

        def _place_column(groups_list, left, width, top, bottom):
            cur_y = top
            for grp in groups_list:
                for elem in grp:
                    avail_h = bottom - cur_y
                    est_h = self._estimate_element_height(elem, width)
                    if avail_h < int(Inches(0.3)) or (est_h > avail_h and avail_h < int(Inches(0.8))):
                        return cur_y
                    max_h = int(max(Inches(0.3), avail_h - Inches(0.1)))
                    h = self._add_element(slide, elem, left, cur_y, width, max_h)
                    if h > 0:
                        cur_y += int(h) + int(Inches(0.08))
            return cur_y

        _place_column(col1_groups, col_left_1, col_width, col_top, slide_bottom)
        _place_column(col2_groups, col_left_2, col_width, col_top, slide_bottom)

    def _create_fallback_slide(self, slide_data: SlideData, error: str):
        """Create a simple fallback slide when normal creation fails."""
        slide = self._add_blank_slide()
        self._set_background(slide)

        # Add title
        self._add_title_shape(slide, slide_data.title)

        # Add error note
        shape = slide.shapes.add_textbox(
            self.MARGIN_LEFT,
            self.MARGIN_TOP,
            int(self.SLIDE_WIDTH - self.MARGIN_LEFT - self.MARGIN_RIGHT),
            Inches(1)
        )
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = f"[Content generation error: {error}]"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(200, 100, 100)

    def _add_element(self, slide, element: ContentElement, left, top, width, max_height=None) -> float:
        """
        Add a content element to the slide.

        Returns:
            Height of the added element in EMUs
        """
        if max_height is None:
            max_height = self.SLIDE_HEIGHT - self.MARGIN_TOP - self.MARGIN_BOTTOM

        if element.type == ContentType.PARAGRAPH:
            return self._add_paragraph_element(slide, element, left, top, width, max_height)
        elif element.type == ContentType.HEADING:
            return self._add_heading_element(slide, element, left, top, width)
        elif element.type in (ContentType.BULLET_LIST, ContentType.NUMBERED_LIST):
            return self._add_list_element(slide, element, left, top, width, max_height)
        elif element.type == ContentType.CODE_BLOCK:
            return self._add_code_element(slide, element, left, top, width, max_height)
        elif element.type == ContentType.TABLE:
            return self._add_table_element(slide, element, left, top, width, max_height)
        elif element.type == ContentType.IMAGE:
            return self._add_image_element(slide, element, left, top, width, max_height)
        else:
            return Inches(0)

    def _add_paragraph_element(self, slide, element: ContentElement, left, top, width, max_height=None) -> float:
        """Add paragraph text with dynamic height."""
        text = element.content
        font_size = self.theme.body_size
        # Estimate lines: ~70 chars per line at body_size for standard width
        chars_per_line = max(40, int(width / Inches(1) * 5))
        est_lines = max(1, len(text) // chars_per_line + (1 if len(text) % chars_per_line else 0))
        height = Inches(max(0.35, est_lines * 0.3))
        if max_height:
            height = min(height, max_height)

        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(*self.theme.text_color)
        p.font.name = self.theme.body_font
        p.space_after = Pt(2)

        return height

    def _add_heading_element(self, slide, element: ContentElement, left, top, width) -> float:
        """Add sub-heading text."""
        level = element.metadata.get("level", 3)
        size = max(14, self.theme.body_size + (6 - level) * 2)

        height = Inches(0.4)
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame

        p = tf.paragraphs[0]
        p.text = element.content
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.theme.title_color)
        p.font.name = self.theme.title_font

        return height

    def _add_list_element(self, slide, element: ContentElement, left, top, width, max_height=None) -> float:
        """Add bullet or numbered list with dynamic height."""
        items = element.content
        is_numbered = element.type == ContentType.NUMBERED_LIST

        line_height = Inches(0.32)
        # Account for long items that wrap
        total_lines = 0
        chars_per_line = max(40, int(width / Inches(1) * 5))
        for item in items:
            text_len = len(item['text']) + item.get("indent", 0) * 4 + 3  # prefix + indent
            total_lines += max(1, text_len // chars_per_line + (1 if text_len % chars_per_line else 0))

        height = line_height * total_lines
        if max_height:
            height = min(height, max_height - Inches(0.1))

        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            indent = item.get("indent", 0)
            prefix = f"{i + 1}. " if is_numbered else "• "
            indent_str = "    " * indent

            p.text = f"{indent_str}{prefix}{item['text']}"
            p.font.size = Pt(self.theme.body_size)
            p.font.color.rgb = RGBColor(*self.theme.text_color)
            p.font.name = self.theme.body_font
            p.space_after = Pt(4)

        return height

    def _add_code_element(self, slide, element: ContentElement, left, top, width, max_height=None) -> float:
        """Add code block with simple styling."""
        code = element.content
        lines = code.split('\n')
        line_count = min(len(lines), 15)  # Limit lines

        line_height = Inches(0.25)
        height = line_height * line_count + Inches(0.3)
        if max_height:
            height = min(height, max_height)

        # Background shape
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(*self.theme.code_bg_color)
        bg_shape.line.fill.background()

        # Code text
        text_shape = slide.shapes.add_textbox(
            left + Inches(0.1),
            top + Inches(0.1),
            width - Inches(0.2),
            height - Inches(0.2)
        )
        tf = text_shape.text_frame
        tf.word_wrap = False

        display_code = '\n'.join(lines[:line_count])
        if len(lines) > line_count:
            display_code += f"\n... ({len(lines) - line_count} more lines)"

        p = tf.paragraphs[0]
        p.text = display_code
        p.font.size = Pt(self.theme.code_size)
        p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(*self.theme.text_color)

        return height

    def _add_table_element(self, slide, element: ContentElement, left, top, width, max_height=None) -> float:
        """Add table."""
        data = element.content
        headers = data.get("headers", [])
        rows = data.get("rows", [])

        if not headers and not rows:
            return Inches(0)

        col_count = len(headers) if headers else len(rows[0]) if rows else 0
        row_count = len(rows) + (1 if headers else 0)

        if col_count == 0 or row_count == 0:
            return Inches(0)

        row_height = Inches(0.35)
        height = row_height * row_count
        if max_height:
            height = min(height, max_height)
            row_count_fit = max(1, int(height / row_height))

        table = slide.shapes.add_table(
            row_count, col_count,
            left, top, width, height
        ).table

        # Style table
        for i, cell in enumerate(table.rows[0].cells if headers else []):
            cell.text = headers[i] if i < len(headers) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*self.theme.accent_color)
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(self.theme.body_size - 2)
            p.font.color.rgb = RGBColor(255, 255, 255)

        start_row = 1 if headers else 0
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_text in enumerate(row_data):
                if col_idx < col_count:
                    cell = table.rows[start_row + row_idx].cells[col_idx]
                    cell.text = str(cell_text)
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(self.theme.body_size - 2)
                    p.font.color.rgb = RGBColor(*self.theme.text_color)

        return height

    def _add_image_element(self, slide, element: ContentElement, left, top, width, max_height=None) -> float:
        """Add image from local file with optional position control."""
        img_path = element.content
        meta = element.metadata or {}

        if not os.path.exists(img_path):
            self.warnings.append(f"Image not found: {img_path}")
            return Inches(0)

        try:
            # Check for position preset
            position = meta.get("position", "")

            # Parse custom width/height from metadata using actual slide dimensions
            slide_w, slide_h = self._get_slide_dimensions()
            custom_width = self._parse_dimension(meta.get("width", ""), slide_w)
            custom_height = self._parse_dimension(meta.get("height", ""), slide_h)

            # Determine target width for aspect ratio calculation
            target_width = custom_width if custom_width else width
            img_max_height = custom_height if custom_height else (max_height if max_height else Inches(4))

            if HAS_PIL:
                with Image.open(img_path) as img:
                    img_width, img_height = img.size
                    aspect = img_width / img_height

                    if target_width / aspect > img_max_height:
                        final_height = img_max_height
                        final_width = img_max_height * aspect
                    else:
                        final_width = target_width
                        final_height = target_width / aspect
            else:
                final_width = target_width
                final_height = img_max_height

            # Round to integer EMU values (python-pptx requires int)
            final_width = int(round(final_width))
            final_height = int(round(final_height))

            # Calculate position based on preset or flow position
            if position:
                img_left, img_top = self._resolve_image_position(
                    position, final_width, final_height, meta
                )
            else:
                # Default: center horizontally in the flow
                img_left = left + (width - final_width) / 2
                img_top = top

            img_left = int(round(img_left))
            img_top = int(round(img_top))

            slide.shapes.add_picture(img_path, img_left, img_top, final_width, final_height)

            # Positioned images don't consume flow space (layout handled by _calc_content_area)
            if position:
                return Inches(0)
            return final_height

        except Exception as e:
            self.warnings.append(f"Failed to add image {img_path}: {str(e)}")
            return Inches(0)

    def _parse_dimension(self, val_str: str, total_size) -> float:
        """Parse dimension string like '50%', '3in', '200px' to EMU."""
        if not val_str:
            return 0
        val_str = val_str.strip()
        if val_str.endswith('%'):
            try:
                pct = float(val_str[:-1]) / 100
                return total_size * pct
            except ValueError:
                return 0
        elif val_str.endswith('in'):
            try:
                return Inches(float(val_str[:-2]))
            except ValueError:
                return 0
        elif val_str.endswith('px'):
            try:
                return Inches(float(val_str[:-2]) / 96)
            except ValueError:
                return 0
        else:
            try:
                return Inches(float(val_str))
            except ValueError:
                return 0

    def _resolve_image_position(self, position: str, img_width, img_height, meta: dict, title_top_offset: float = None):
        """Resolve image position preset to (left, top) coordinates.

        Args:
            title_top_offset: Bottom of title area in EMU. Top-positioned images
                are placed below this to avoid overlap.
        """
        margin = Inches(0.3)
        slide_w, slide_h = self._get_slide_dimensions()

        # Default top for top-positioned images: below title or 1.2in
        top_y = title_top_offset if title_top_offset is not None else Inches(1.2)

        # Parse custom x/y if provided
        custom_x = self._parse_dimension(meta.get("x", ""), slide_w)
        custom_y = self._parse_dimension(meta.get("y", ""), slide_h)

        positions = {
            "top-left": (margin, top_y),
            "top-center": ((slide_w - img_width) / 2, top_y),
            "top-right": (slide_w - img_width - margin, top_y),
            "center-left": (margin, (slide_h - img_height) / 2),
            "center": ((slide_w - img_width) / 2, max(top_y, (slide_h - img_height) / 2)),
            "center-right": (slide_w - img_width - margin, (slide_h - img_height) / 2),
            "bottom-left": (margin, slide_h - img_height - margin),
            "bottom-center": ((slide_w - img_width) / 2, slide_h - img_height - margin),
            "bottom-right": (slide_w - img_width - margin, slide_h - img_height - margin),
        }

        if custom_x or custom_y:
            return (custom_x if custom_x else margin, custom_y if custom_y else top_y)

        return positions.get(position, positions["center"])

    def _add_chart_from_table(self, slide, element: ContentElement, left, top, width, height):
        """Create a chart from table data."""
        data = element.content
        headers = data.get("headers", [])
        rows = data.get("rows", [])

        if len(headers) < 2 or not rows:
            # Fall back to table
            self._add_table_element(slide, element, left, top, width)
            return

        try:
            chart_data = CategoryChartData()
            chart_data.categories = [row[0] for row in rows]

            # Add series for each numeric column
            for col_idx in range(1, len(headers)):
                values = []
                for row in rows:
                    if col_idx < len(row):
                        try:
                            values.append(float(row[col_idx].replace(',', '')))
                        except (ValueError, AttributeError):
                            values.append(0)
                    else:
                        values.append(0)
                chart_data.add_series(headers[col_idx], values)

            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                left, top, width, height,
                chart_data
            ).chart

            # Style chart
            chart.has_legend = len(headers) > 2

        except Exception as e:
            self.warnings.append(f"Chart creation failed: {str(e)}")
            self._add_table_element(slide, element, left, top, width)

    def _add_chart_element(self, slide, element: ContentElement, left, top, width, height):
        """Add chart from chart data."""
        # Placeholder for direct chart data support
        pass

    def get_warnings(self) -> List[str]:
        """Return warnings generated during generation."""
        return self.warnings


def generate_pptx(
    slides: List[SlideData],
    output_path: str,
    theme: str = "business",
    template_path: Optional[str] = None,
    progress_callback=None
) -> Tuple[bool, List[str]]:
    """
    Convenience function to generate PPTX.

    Args:
        slides: List of SlideData objects
        output_path: Output file path
        theme: Theme name
        template_path: Optional custom template path
        progress_callback: Optional callback(current, total, message)

    Returns:
        Tuple of (success, warnings)
    """
    generator = PPTXGenerator(theme=theme, template_path=template_path)
    if progress_callback:
        generator.set_progress_callback(progress_callback)

    success = generator.generate(slides, output_path)
    return success, generator.get_warnings()


if __name__ == "__main__":
    # Test with sample data
    from md_parser import parse_markdown

    sample_md = """
# Test Presentation

A demo presentation

## Introduction

Welcome to this test presentation.

- First point
- Second point
- Third point

## Code Example

```python
def hello():
    print("Hello!")
```

## Data

| Item | Value |
|------|-------|
| A    | 100   |
| B    | 200   |
| C    | 150   |

## Thank You

Questions?
"""

    slides, parse_warnings = parse_markdown(sample_md)
    print(f"Parsed {len(slides)} slides")

    def progress(current, total, msg):
        print(f"[{current}/{total}] {msg}")

    success, gen_warnings = generate_pptx(
        slides,
        "test_output.pptx",
        theme="business",
        progress_callback=progress
    )

    print(f"\nGeneration {'succeeded' if success else 'failed'}")
    if parse_warnings:
        print(f"Parse warnings: {parse_warnings}")
    if gen_warnings:
        print(f"Generation warnings: {gen_warnings}")
